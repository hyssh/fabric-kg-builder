"""PPTX source adapter (EXT-006).

Extracts from PowerPoint presentations:
- Slide order (slide_number, sort_order)
- Slide notes (element_type="notes")
- Text from shapes (element_type="paragraph")
- Tables from table shapes (element_type="table" / "table_row")
- Image shape references (element_type="image_ref")
- Hyperlinks with anchor/target/position (EXT-003)

Limits (EXT-009)
----------------
- File size: ``adapter.MAX_FILE_BYTES`` (200 MB).
- Slide count: ``adapter.MAX_SLIDES`` (1 000).
"""

from __future__ import annotations

import hashlib
import html as html_escape_mod
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from fabric_kg_builder.model.ids import (
    content_hash as compute_content_hash,
    make_document_element_id,
    make_source_file_id,
)
from fabric_kg_builder.model.schemas import DocumentElementRow, SourceFileRow

from .adapter import (
    ADAPTER_CONTRACT_VERSION,
    MAX_FILE_BYTES,
    MAX_SLIDES,
    AdapterError,
    AdapterResult,
    FailureType,
    HyperlinkRecord,
    check_file_size,
)
from .media_type import validate_extension_vs_signature

_ADAPTER_NAME = "pptx_extractor"
_ADAPTER_VERSION = "1.0.0"


def _file_hash(path: Path) -> str:
    """SHA-256 via streaming read."""
    h = hashlib.sha256()
    with open(path, "rb") as fh:
        for chunk in iter(lambda: fh.read(65536), b""):
            h.update(chunk)
    return h.hexdigest()


def _canonical_path(path: Path, project_root: Path | None) -> str:
    root = project_root or Path.cwd()
    try:
        return path.resolve().relative_to(root.resolve()).as_posix()
    except ValueError:
        return path.as_posix()


def _shape_hyperlinks(
    shape: Any,
    source_element_id: str,
    slide_number: int,
    shape_name: str = "",
) -> list[HyperlinkRecord]:
    """Extract hyperlinks from a pptx shape's text runs with exact locators."""
    links: list[HyperlinkRecord] = []
    try:
        has_tf = shape.has_text_frame
    except AttributeError:
        return links
    if not has_tf:
        return links
    for para_idx, para in enumerate(shape.text_frame.paragraphs):
        for run_idx, run in enumerate(para.runs):
            try:
                hl = run.hyperlink
                if hl is None:
                    continue
                address = getattr(hl, "address", None) or ""
                if address:
                    locator = json.dumps({
                        "type": "pptx",
                        "slide": slide_number,
                        "shape_name": shape_name,
                        "paragraph_index": para_idx,
                        "run_index": run_idx,
                    })
                    links.append(
                        HyperlinkRecord(
                            anchor=run.text or "",
                            target=address,
                            source_element_id=source_element_id,
                            page_number=slide_number,
                            source_locator_json=locator,
                        )
                    )
            except AttributeError:
                continue
    return links


def _table_to_html(table: Any) -> str:
    """Convert a pptx table to HTML string."""
    rows_html: list[str] = []
    for row in table.rows:
        cells_html = "".join(
            f"<td>{html_escape_mod.escape(cell.text_frame.text.strip())}</td>"
            for cell in row.cells
        )
        rows_html.append(f"<tr>{cells_html}</tr>")
    return f"<table>{''.join(rows_html)}</table>"


class PptxExtractor:
    """Source adapter for PowerPoint PPTX files (EXT-006)."""

    @staticmethod
    def extract(
        path: str | Path,
        *,
        project_root: str | Path | None = None,
        validate_mime: bool = True,
        max_slides: int = MAX_SLIDES,
        max_bytes: int = MAX_FILE_BYTES,
    ) -> AdapterResult:
        """Extract document elements from a PPTX file.

        Parameters
        ----------
        path:
            Path to the ``.pptx`` file.
        project_root:
            Optional root for canonical relative-path computation.
        validate_mime:
            When True (default), validate magic bytes match ``.pptx`` extension.
        max_slides:
            Reject files exceeding this slide count (EXT-009).
        max_bytes:
            Reject files exceeding this byte size (EXT-009).

        Returns
        -------
        AdapterResult
            ``adapter_name="pptx_extractor"``,
            ``page_count=len(slides)``,
            ``document_elements`` contains one element per slide/text/table/image,
            ``hyperlinks`` carries all hyperlinks found across slides.
        """
        path = Path(path)
        if not path.exists():
            raise AdapterError(
                FailureType.NOT_FOUND,
                f"Source file not found: {path}",
                source_locator=str(path),
            )

        check_file_size(path, max_bytes)

        detected_mime = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        if validate_mime:
            detected_mime = validate_extension_vs_signature(path)

        now = datetime.now(timezone.utc)
        file_hash = _file_hash(path)
        can_path = _canonical_path(path, Path(project_root) if project_root else None)
        source_file_id = make_source_file_id(can_path, file_hash)

        try:
            from pptx import Presentation  # noqa: PLC0415
            from pptx.util import Pt  # noqa: PLC0415 (for type reference only)
            from pptx.enum.shapes import MSO_SHAPE_TYPE  # noqa: PLC0415
        except ImportError as exc:
            raise ImportError(
                "python-pptx is required to use the PPTX adapter. "
                "Install it with: pip install python-pptx"
            ) from exc

        try:
            prs = Presentation(path)
        except Exception as exc:
            raise AdapterError(
                FailureType.CORRUPT,
                f"Cannot open PPTX '{path.name}': {exc}",
                source_locator=str(path),
            ) from exc

        slides = prs.slides
        num_slides = len(slides)

        if num_slides > max_slides:
            raise AdapterError(
                FailureType.TOO_MANY_SLIDES,
                f"Presentation '{path.name}' has {num_slides} slides, which "
                f"exceeds the {max_slides}-slide limit.",
                source_locator=str(path),
            )

        source_file = SourceFileRow(
            source_file_id=source_file_id,
            path=can_path,
            filename=path.name,
            source_type="pptx",
            content_hash=file_hash,
            byte_size=path.stat().st_size,
            ingested_at=now,
            row_count=num_slides,
            notes=f"slides={num_slides}",
        )

        document_elements: list[DocumentElementRow] = []
        hyperlinks: list[HyperlinkRecord] = []
        sort_order = 0

        for slide_idx, slide in enumerate(slides):
            slide_number = slide_idx + 1

            # Slide element (parent for all content on this slide)
            slide_title = ""
            try:
                title_shape = slide.shapes.title
                if title_shape is not None:
                    slide_title = title_shape.text.strip()
            except AttributeError:
                pass  # Some slide layouts have no title placeholder

            slide_content = slide_title or f"Slide {slide_number}"
            slide_hash = compute_content_hash(f"slide:{slide_number}:{slide_content}")
            slide_elem_id = make_document_element_id(
                source_file_id, "section", slide_number, sort_order, slide_hash
            )
            document_elements.append(
                DocumentElementRow(
                    document_element_id=slide_elem_id,
                    source_file_id=source_file_id,
                    element_type="section",
                    title=slide_content,
                    content=slide_content,
                    page_number=slide_number,
                    sort_order=sort_order,
                    content_hash=slide_hash,
                    extracted_at=now,
                )
            )
            sort_order += 1

            # Notes
            notes_text = ""
            try:
                if slide.has_notes_slide:
                    notes_tf = slide.notes_slide.notes_text_frame
                    notes_text = notes_tf.text.strip()
            except AttributeError:
                pass  # notes_slide / notes_text_frame may not exist on all slide types

            if notes_text:
                notes_hash = compute_content_hash(notes_text)
                notes_elem_id = make_document_element_id(
                    source_file_id, "notes", slide_number, sort_order, notes_hash
                )
                document_elements.append(
                    DocumentElementRow(
                        document_element_id=notes_elem_id,
                        source_file_id=source_file_id,
                        element_type="notes",
                        parent_element_id=slide_elem_id,
                        content=notes_text,
                        page_number=slide_number,
                        sort_order=sort_order,
                        content_hash=notes_hash,
                        extracted_at=now,
                    )
                )
                sort_order += 1

            # Shapes
            for shape in slide.shapes:
                shape_type = shape.shape_type

                # Picture / image reference
                if shape_type == MSO_SHAPE_TYPE.PICTURE:
                    img_name = getattr(shape, "name", f"image_{sort_order}")
                    img_content = f"{path.name}#slide{slide_number}/{img_name}"
                    img_hash = compute_content_hash(img_content)
                    img_elem_id = make_document_element_id(
                        source_file_id, "image_ref", slide_number, sort_order, img_hash
                    )
                    document_elements.append(
                        DocumentElementRow(
                            document_element_id=img_elem_id,
                            source_file_id=source_file_id,
                            element_type="image_ref",
                            parent_element_id=slide_elem_id,
                            title=img_name,
                            content=img_content,
                            page_number=slide_number,
                            sort_order=sort_order,
                            content_hash=img_hash,
                            extracted_at=now,
                        )
                    )
                    sort_order += 1
                    continue

                # Table shape
                if shape.has_table:
                    tbl = shape.table
                    table_html = _table_to_html(tbl)
                    table_text = " | ".join(
                        " ".join(
                            cell.text_frame.text.strip()
                            for cell in row.cells
                        ).strip()
                        for row in tbl.rows
                    ).strip()
                    table_hash = compute_content_hash(table_text)
                    table_elem_id = make_document_element_id(
                        source_file_id, "table", slide_number, sort_order, table_hash
                    )
                    document_elements.append(
                        DocumentElementRow(
                            document_element_id=table_elem_id,
                            source_file_id=source_file_id,
                            element_type="table",
                            parent_element_id=slide_elem_id,
                            content=table_text,
                            content_html=table_html,
                            page_number=slide_number,
                            sort_order=sort_order,
                            content_hash=table_hash,
                            extracted_at=now,
                        )
                    )
                    sort_order += 1

                    for row_idx, row in enumerate(tbl.rows):
                        row_text = " | ".join(
                            cell.text_frame.text.strip()
                            for cell in row.cells
                        ).strip()
                        if not row_text:
                            continue
                        row_html = (
                            "<tr>"
                            + "".join(
                                f"<td>{html_escape_mod.escape(cell.text_frame.text.strip())}</td>"
                                for cell in row.cells
                            )
                            + "</tr>"
                        )
                        row_hash = compute_content_hash(row_text)
                        document_elements.append(
                            DocumentElementRow(
                                document_element_id=make_document_element_id(
                                    source_file_id, "table_row", slide_number, sort_order, row_hash
                                ),
                                source_file_id=source_file_id,
                                element_type="table_row",
                                parent_element_id=table_elem_id,
                                content=row_text,
                                content_html=row_html,
                                page_number=slide_number,
                                sort_order=sort_order,
                                row_index=row_idx,
                                content_hash=row_hash,
                                extracted_at=now,
                            )
                        )
                        sort_order += 1
                    continue

                # Text frame shapes (skip title, already captured)
                if shape.has_text_frame:
                    # Collect hyperlinks from this shape's text runs
                    para_elem_id_for_links: str | None = None

                    for para in shape.text_frame.paragraphs:
                        para_text = para.text.strip()
                        if not para_text:
                            continue
                        para_hash = compute_content_hash(para_text)
                        para_elem_id = make_document_element_id(
                            source_file_id, "paragraph", slide_number, sort_order, para_hash
                        )
                        document_elements.append(
                            DocumentElementRow(
                                document_element_id=para_elem_id,
                                source_file_id=source_file_id,
                                element_type="paragraph",
                                parent_element_id=slide_elem_id,
                                content=para_text,
                                page_number=slide_number,
                                sort_order=sort_order,
                                content_hash=para_hash,
                                extracted_at=now,
                            )
                        )
                        para_elem_id_for_links = para_elem_id
                        sort_order += 1

                    # Extract hyperlinks from runs
                    shape_name_str = getattr(shape, "name", "") or ""
                    links = _shape_hyperlinks(shape, para_elem_id_for_links or slide_elem_id, slide_number, shape_name=shape_name_str)
                    hyperlinks.extend(links)

        return AdapterResult(
            adapter_name=_ADAPTER_NAME,
            adapter_version=_ADAPTER_VERSION,
            detected_media_type=detected_mime,
            source_locator=f"file://{path.resolve().as_posix()}",
            source_file=source_file,
            document_elements=document_elements,
            page_count=num_slides,
            hyperlinks=hyperlinks,
        )
